"""
Document processing module for the multimodal RAG chatbot
"""

import os
import sys
import logging
from typing import Dict, List, Any, Optional, Tuple
from datetime import datetime
import tempfile

# Add the project root to the path
sys.path.append(os.path.dirname(os.path.dirname(os.path.abspath(__file__))))

from config import config
from utils.base_utils import chunk_text, get_file_extension, get_file_type
from utils.db_manager import get_db_manager

# Import document processing libraries
import PyPDF2
from docx import Document
from pptx import Presentation
import openpyxl
import pytesseract
from pdf2image import convert_from_bytes
from PIL import Image

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Base class for document processing"""
    
    def __init__(self):
        """Initialize document processor"""
        self.db_manager = get_db_manager()
    
    def process_document(self, file_path: str, file_content: bytes, 
                         metadata: Dict[str, Any]) -> str:
        """
        Process a document file
        
        Args:
            file_path: Path to the document file
            file_content: Binary content of the file
            metadata: Additional metadata for the document
            
        Returns:
            Document ID
        """
        try:
            # Extract text based on file type
            file_extension = get_file_extension(file_path)
            text = self._extract_text(file_content, file_extension)
            
            # Store file in GridFS
            file_id = self.db_manager.store_file(
                os.path.basename(file_path),
                file_content,
                metadata
            )
            
            # Create document metadata
            document_data = {
                "filename": os.path.basename(file_path),
                "file_type": get_file_type(file_path),
                "file_extension": file_extension,
                "file_size": len(file_content),
                "file_id": file_id,
                "text_length": len(text) if text else 0,
                "created_at": datetime.now(),
                "metadata": metadata
            }
            
            # Store document metadata
            document_id = self.db_manager.store_document_metadata(document_data)
            
            # Process text chunks
            if text:
                self._process_chunks(text, document_id)
            
            return document_id
        
        except Exception as e:
            logger.error(f"Error processing document {file_path}: {str(e)}")
            raise
    
    def _extract_text(self, file_content: bytes, file_extension: str) -> str:
        """
        Extract text from document based on file type
        
        Args:
            file_content: Binary content of the file
            file_extension: File extension
            
        Returns:
            Extracted text
        """
        if file_extension == ".pdf":
            return self._extract_text_from_pdf(file_content)
        elif file_extension == ".docx":
            return self._extract_text_from_docx(file_content)
        elif file_extension == ".pptx":
            return self._extract_text_from_pptx(file_content)
        elif file_extension == ".xlsx":
            return self._extract_text_from_xlsx(file_content)
        elif file_extension == ".txt" or file_extension == ".md":
            return file_content.decode('utf-8', errors='replace')
        elif file_extension == ".csv":
            return file_content.decode('utf-8', errors='replace')
        elif file_extension == ".json":
            return file_content.decode('utf-8', errors='replace')
        elif file_extension == ".html":
            return self._extract_text_from_html(file_content)
        else:
            logger.warning(f"Unsupported document type: {file_extension}")
            return ""
    
    def _extract_text_from_pdf(self, file_content: bytes) -> str:
        """Extract text from PDF file"""
        try:
            text = ""
            with tempfile.NamedTemporaryFile(delete=False) as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Try text extraction first
            try:
                with open(temp_file_path, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    for page_num in range(len(pdf_reader.pages)):
                        page = pdf_reader.pages[page_num]
                        text += page.extract_text() + "\n\n"
            except Exception as e:
                logger.warning(f"Error extracting text from PDF: {str(e)}")
            
            # If text extraction failed or returned little text, try OCR
            if len(text.strip()) < 100:
                logger.info("PDF text extraction yielded little text, trying OCR")
                ocr_text = self._ocr_pdf(temp_file_path)
                if ocr_text:
                    text = ocr_text
            
            # Clean up temp file
            os.unlink(temp_file_path)
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PDF: {str(e)}")
            return ""
    
    def _ocr_pdf(self, pdf_path: str) -> str:
        """Perform OCR on PDF file"""
        try:
            images = convert_from_bytes(open(pdf_path, 'rb').read())
            text = ""
            
            for i, image in enumerate(images):
                text += pytesseract.image_to_string(image, lang='fra+eng') + "\n\n"
            
            return text
        except Exception as e:
            logger.error(f"Error performing OCR on PDF: {str(e)}")
            return ""
    
    def _extract_text_from_docx(self, file_content: bytes) -> str:
        """Extract text from DOCX file"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.docx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            doc = Document(temp_file_path)
            text = "\n".join([paragraph.text for paragraph in doc.paragraphs])
            
            # Clean up temp file
            os.unlink(temp_file_path)
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from DOCX: {str(e)}")
            return ""
    
    def _extract_text_from_pptx(self, file_content: bytes) -> str:
        """Extract text from PPTX file"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.pptx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            prs = Presentation(temp_file_path)
            text = ""
            
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
                text += "\n"
            
            # Clean up temp file
            os.unlink(temp_file_path)
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from PPTX: {str(e)}")
            return ""
    
    def _extract_text_from_xlsx(self, file_content: bytes) -> str:
        """Extract text from XLSX file"""
        try:
            with tempfile.NamedTemporaryFile(delete=False, suffix='.xlsx') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            workbook = openpyxl.load_workbook(temp_file_path)
            text = ""
            
            for sheet_name in workbook.sheetnames:
                sheet = workbook[sheet_name]
                text += f"Sheet: {sheet_name}\n"
                
                for row in sheet.iter_rows(values_only=True):
                    row_text = " | ".join([str(cell) if cell is not None else "" for cell in row])
                    if row_text.strip():
                        text += row_text + "\n"
                
                text += "\n"
            
            # Clean up temp file
            os.unlink(temp_file_path)
            
            return text
        except Exception as e:
            logger.error(f"Error extracting text from XLSX: {str(e)}")
            return ""
    
    def _extract_text_from_html(self, file_content: bytes) -> str:
        """Extract text from HTML file"""
        try:
            from bs4 import BeautifulSoup
            soup = BeautifulSoup(file_content, 'html.parser')
            return soup.get_text(separator='\n')
        except ImportError:
            logger.warning("BeautifulSoup not installed, returning raw HTML")
            return file_content.decode('utf-8', errors='replace')
        except Exception as e:
            logger.error(f"Error extracting text from HTML: {str(e)}")
            return ""
    
    def _process_chunks(self, text: str, document_id: str) -> List[str]:
        """
        Process text into chunks and store in database
        
        Args:
            text: Document text
            document_id: Document ID
            
        Returns:
            List of chunk IDs
        """
        # Split text into chunks
        text_chunks = chunk_text(text, config.CHUNK_SIZE, config.CHUNK_OVERLAP)
        
        # Prepare chunks data
        chunks_data = []
        for i, chunk_text in enumerate(text_chunks):
            chunks_data.append({
                "document_id": document_id,
                "chunk_index": i,
                "text": chunk_text,
                "char_count": len(chunk_text),
                "created_at": datetime.now()
            })
        
        # Store chunks
        if chunks_data:
            return self.db_manager.store_chunks(chunks_data)
        
        return []


class ImageProcessor:
    """Class for processing image files"""
    
    def __init__(self):
        """Initialize image processor"""
        self.db_manager = get_db_manager()
    
    def process_image(self, file_path: str, file_content: bytes, 
                      metadata: Dict[str, Any]) -> str:
        """
        Process an image file
        
        Args:
            file_path: Path to the image file
            file_content: Binary content of the file
            metadata: Additional metadata for the image
            
        Returns:
            Image document ID
        """
        try:
            # Extract text from image using OCR
            text = self._extract_text_from_image(file_content)
            
            # Store file in GridFS
            file_id = self.db_manager.store_file(
                os.path.basename(file_path),
                file_content,
                metadata
            )
            
            # Create image metadata
            image_data = {
                "filename": os.path.basename(file_path),
                "file_type": "image",
                "file_extension": get_file_extension(file_path),
                "file_size": len(file_content),
                "file_id": file_id,
                "text": text,
                "text_length": len(text) if text else 0,
                "created_at": datetime.now(),
                "metadata": metadata
            }
            
            # Store image metadata
            image_id = self.db_manager.images.insert_one(image_data).inserted_id
            
            return str(image_id)
        
        except Exception as e:
            logger.error(f"Error processing image {file_path}: {str(e)}")
            raise
    
    def _extract_text_from_image(self, file_content: bytes) -> str:
        """
        Extract text from image using OCR
        
        Args:
            file_content: Binary content of the image
            
        Returns:
            Extracted text
        """
        try:
            # Create a temporary file
            with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_file:
                temp_file.write(file_content)
                temp_file_path = temp_file.name
            
            # Open the image
            image = Image.open(temp_file_path)
            
            # Perform OCR
            text = pytesseract.image_to_string(image, lang='fra+eng')
            
            # Clean up temp file
            os.unlink(temp_file_path)
            
            return text
        
        except Exception as e:
            logger.error(f"Error extracting text from image: {str(e)}")
            return ""


class AudioProcessor:
    """Class for processing audio files"""
    
    def __init__(self):
        """Initialize audio processor"""
        self.db_manager = get_db_manager()
    
    def process_audio(self, file_path: str, file_content: bytes, 
                      metadata: Dict[str, Any]) -> str:
        """
        Process an audio file
        
        Args:
            file_path: Path to the audio file
            file_content: Binary content of the file
            metadata: Additional metadata for the audio
            
        Returns:
            Audio document ID
        """
        try:
            # Store file in GridFS
            file_id = self.db_manager.store_file(
                os.path.basename(file_path),
                file_content,
                metadata
            )
            
            # Create audio metadata
            audio_data = {
                "filename": os.path.basename(file_path),
                "file_type": "audio",
                "file_extension": get_file_extension(file_path),
                "file_size": len(file_content),
                "file_id": file_id,
                "transcription": "",  # Will be updated when transcription is available
                "created_at": datetime.now(),
                "metadata": metadata
            }
            
            # Store audio metadata
            audio_id = self.db_manager.audio.insert_one(audio_data).inserted_id
            
            # Note: Actual transcription would be done asynchronously due to resource constraints
            # For now, we'll just store the audio file and metadata
            
            return str(audio_id)
        
        except Exception as e:
            logger.error(f"Error processing audio {file_path}: {str(e)}")
            raise


class VideoProcessor:
    """Class for processing video files"""
    
    def __init__(self):
        """Initialize video processor"""
        self.db_manager = get_db_manager()
    
    def process_video(self, file_path: str, file_content: bytes, 
                      metadata: Dict[str, Any]) -> str:
        """
        Process a video file
        
        Args:
            file_path: Path to the video file
            file_content: Binary content of the file
            metadata: Additional metadata for the video
            
        Returns:
            Video document ID
        """
        try:
            # Store file in GridFS
            file_id = self.db_manager.store_file(
                os.path.basename(file_path),
                file_content,
                metadata
            )
            
            # Create video metadata
            video_data = {
                "filename": os.path.basename(file_path),
                "file_type": "video",
                "file_extension": get_file_extension(file_path),
                "file_size": len(file_content),
                "file_id": file_id,
                "transcription": "",  # Will be updated when transcription is available
                "created_at": datetime.now(),
                "metadata": metadata
            }
            
            # Store video metadata
            video_id = self.db_manager.video.insert_one(video_data).inserted_id
            
            # Note: Actual transcription and frame extraction would be done asynchronously
            # due to resource constraints
            # For now, we'll just store the video file and metadata
            
            return str(video_id)
        
        except Exception as e:
            logger.error(f"Error processing video {file_path}: {str(e)}")
            raise


# Factory function to get the appropriate processor
def get_processor(file_type: str):
    """
    Get the appropriate processor for a file type
    
    Args:
        file_type: Type of file (document, image, audio, video)
        
    Returns:
        Processor instance
    """
    if file_type == "document":
        return DocumentProcessor()
    elif file_type == "image":
        return ImageProcessor()
    elif file_type == "audio":
        return AudioProcessor()
    elif file_type == "video":
        return VideoProcessor()
    else:
        raise ValueError(f"Unsupported file type: {file_type}")
